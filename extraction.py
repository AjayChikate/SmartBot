import os
import io

from dotenv import load_dotenv

import pdfplumber
from PyPDF2 import PdfReader
from pdf2image import convert_from_bytes
from ocr import ocr_image

from docx import Document

from pptx import Presentation

from bs4 import BeautifulSoup

def extract_pdf_text(pdf_file, enable_ocr=False, poppler_path=None):
    text = ""
    pdf_bytes = pdf_file.read()
    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:   # High level Structured img, tables and text Extractor, But easily fails in non structured cases so i ave a fallback
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except:
        pass
    if not text.strip():
        try:
            pdf_reader = PdfReader(io.BytesIO(pdf_bytes))  #plain text extraction ,Basic low level reader used in unstructered pdfs as fallback
            for page in pdf_reader.pages:
                t = page.extract_text()
                if t:
                    text += t + "\n"
        except:
            pass
    if enable_ocr or not text.strip():
        try:
            poppler_path = os.getenv("POPPLER_PATH")
            images = convert_from_bytes(pdf_bytes, poppler_path=poppler_path)   # From Scanned Images gives text 
            for img in images:
                ocr_text = ocr_image(img)
                if ocr_text:
                    text += ocr_text + "\n"
        except:
            pass
    return text

def extract_docx_text(docx_file):
    text = ""
    try:
        doc = Document(io.BytesIO(docx_file.read()))
        for para in doc.paragraphs:
            if para.text.strip():
                style = para.style.name if para.style else ""
                if 'Heading' in style:
                    text += f"\n## {para.text}\n"
                else:
                    text += para.text + "\n"
    except:
        pass
    return text

def extract_pptx_text(pptx_file):
    text = ""
    try:
        prs = Presentation(io.BytesIO(pptx_file.read()))
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n### Slide {slide_num}\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes.strip():
                        text += f"[Notes]: {notes}\n"
            except:
                pass
    except:
        pass
    return text

def extract_html_text(html_file):
    text = ""
    try:
        html_content = html_file.read().decode('utf-8', errors='ignore')
        soup = BeautifulSoup(html_content, 'html.parser')
        for script in soup(["script", "style"]):
            script.decompose()
        text = soup.get_text(separator='\n', strip=True)
    except:
        pass
    return text

def extract_text_file(txt_file):
    try:
        return txt_file.read().decode('utf-8', errors='ignore')
    except:
        return ""